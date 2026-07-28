import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BG_COLOR = RGBColor(11, 15, 25)         # #0B0F19
    CARD_BG = RGBColor(15, 23, 42)          # #0F172A
    CARD_BORDER = RGBColor(56, 189, 248)    # #38BDF8
    TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8
    CYAN_ACCENT = RGBColor(56, 189, 248)   # #38BDF8
    INDIGO_ACCENT = RGBColor(129, 140, 248) # #818CF8
    EMERALD_ACCENT = RGBColor(52, 211, 153) # #34D399
    AMBER_ACCENT = RGBColor(251, 191, 36)   # #FBBF24
    ROSE_ACCENT = RGBColor(251, 113, 133)   # #FB7185

    logo_path = os.path.abspath(r"c:\Personal\3D\docs\kle_tech_logo.png")
    pdf_assets_dir = os.path.abspath(r"c:\Personal\3D\docs\pdf_extracted_assets")

    img_cover = os.path.join(pdf_assets_dir, "img_page_1_1.png")
    img_intro = os.path.join(pdf_assets_dir, "img_page_3_1.png")
    img_prob = os.path.join(pdf_assets_dir, "img_page_4_1.png")
    img_obj = os.path.join(pdf_assets_dir, "img_page_5_1.png")
    img_scope = os.path.join(pdf_assets_dir, "img_page_6_1.png")

    def add_blank_slide(bg_color=BG_COLOR):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        return slide

    def add_header(slide, title_text, slide_num):
        header_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        header_shape.line.color.rgb = RGBColor(30, 41, 59)
        
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(1.0), Inches(0.48), height=Inches(0.64))
        
        tf = header_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(1.0)
        p = tf.paragraphs[0]
        p.text = f"  KLE Technological University — Dept. of MCA"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"  {title_text}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = CYAN_ACCENT

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.55), Inches(1.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
        badge.line.color.rgb = CYAN_ACCENT
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = f"Slide {slide_num} / 10"
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 0: TITLE COVER SLIDE WITH IMAGE
    # -------------------------------------------------------------------------
    s0 = add_blank_slide()
    card0 = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    card0.fill.solid()
    card0.fill.fore_color.rgb = CARD_BG
    card0.line.color.rgb = CARD_BORDER

    if os.path.exists(logo_path):
        s0.shapes.add_picture(logo_path, Inches(11.0), Inches(1.1), height=Inches(0.9))

    if os.path.exists(img_cover):
        s0.shapes.add_picture(img_cover, Inches(7.6), Inches(2.2), height=Inches(3.8))

    tf0 = card0.text_frame
    tf0.word_wrap = True
    tf0.margin_left = Inches(0.5)
    tf0.margin_top = Inches(0.4)

    p0 = tf0.paragraphs[0]
    p0.text = "MCA IV SEMESTER MAIN PROJECT PRESENTATION"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = CYAN_ACCENT

    p1 = tf0.add_paragraph()
    p1.text = "Automated Single-Image to 3D Asset &\nPoint Cloud Generation System"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_before = Pt(10)

    p2 = tf0.add_paragraph()
    p2.text = "An end-to-end web system converting single 2D RGB photographs into textured 3D polygon meshes (.GLB) and surface point clouds (.PLY) in ~19.6 seconds."
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)

    # Info Cards
    box1 = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.5), Inches(3.0), Inches(1.6))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(2, 6, 23)
    box1.line.color.rgb = RGBColor(30, 41, 59)
    tf_b1 = box1.text_frame
    p = tf_b1.paragraphs[0]
    p.text = "SUBMITTED BY"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p = tf_b1.add_paragraph()
    p.text = "Rajashekhar B Durgad"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_b1.add_paragraph()
    p.text = "USN: 01FE24MCA027"
    p.font.size = Pt(10)
    p.font.color.rgb = CYAN_ACCENT

    box2 = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.4), Inches(4.5), Inches(3.0), Inches(1.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(2, 6, 23)
    box2.line.color.rgb = RGBColor(30, 41, 59)
    tf_b2 = box2.text_frame
    p = tf_b2.paragraphs[0]
    p.text = "UNDER GUIDANCE OF"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p = tf_b2.add_paragraph()
    p.text = "Prof. Akash Hulkund"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_b2.add_paragraph()
    p.text = "Dept. of MCA, KLE Tech"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------------------
    # SLIDE 1: INTRODUCTION WITH IMAGE
    # -------------------------------------------------------------------------
    s1 = add_blank_slide()
    add_header(s1, "1. Introduction — Overview & Multi-Stage Architecture", 1)

    cards = [
        ("Automated 2D to 3D Pipeline", "Transforms raw 2D photographs into textured 3D polygon meshes (.GLB) and surface point clouds (.PLY) without multi-view setups or manual modeling.", CYAN_ACCENT),
        ("Multi-Stage AI Engine", "Integrates Meta SAM 2.1 (foreground isolation), Tencent Hunyuan3D-2 (triplane synthesis), and Open3D (Poisson disk surface sampling).", INDIGO_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(cards):
        c_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5 + i*2.25), Inches(7.0), Inches(2.1))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = CARD_BG
        c_box.line.color.rgb = color
        tf = c_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    if os.path.exists(img_intro):
        s1.shapes.add_picture(img_intro, Inches(8.1), Inches(1.5), height=Inches(4.35))

    hl = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    hl.fill.solid()
    hl.fill.fore_color.rgb = RGBColor(15, 23, 42)
    hl.line.color.rgb = EMERALD_ACCENT
    tf_hl = hl.text_frame
    p = tf_hl.paragraphs[0]
    p.text = "Key Highlight: Complete processing finished in ~19.6 seconds with a peak GPU memory footprint of 11.4 GB VRAM."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT WITH IMAGE
    # -------------------------------------------------------------------------
    s2 = add_blank_slide()
    add_header(s2, "2. Problem Statement — Challenges & Technical Bottlenecks", 2)

    probs = [
        ("High Manual Overhead in 3D Modeling", "Manual CAD software (Blender, Maya) requires specialized artistic skills and hours to days of labor per single 3D asset.", ROSE_ACCENT),
        ("Photogrammetry & Text-to-3D Limits", "Photogrammetry requires 50+ calibrated camera shots. Text-to-3D generative models lack spatial precision for exact physical objects.", AMBER_ACCENT),
        ("GPU Hardware & VRAM Memory Bottlenecks", "Simultaneous execution of multiple neural networks causes Out-Of-Memory (OOM) crashes on standard workstation GPUs.", INDIGO_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(probs):
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5 + i*1.45), Inches(7.0), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)

    if os.path.exists(img_prob):
        s2.shapes.add_picture(img_prob, Inches(8.1), Inches(1.5), height=Inches(4.2))

    g_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    g_box.line.color.rgb = CYAN_ACCENT
    tf_g = g_box.text_frame
    p = tf_g.paragraphs[0]
    p.text = "The Core Goal: Provide a prompt-free, single-image 3D synthesis web system operating reliably under a 16 GB GPU VRAM ceiling."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 3: OBJECTIVES WITH IMAGE
    # -------------------------------------------------------------------------
    s3 = add_blank_slide()
    add_header(s3, "3. Objectives — System Technical Objectives", 3)

    objs = [
        "Prompt-Free 2D-to-3D Pipeline: Automate single RGB image ingestion without text prompts.",
        "Zero-Shot Background Isolation: Strip background clutter using Meta SAM 2.1 for clean RGBA cutouts.",
        "Rapid 3D Mesh Synthesis: Generate textured watertight .GLB polygon meshes in ~15s using Hunyuan3D-2.",
        "Point Cloud Extraction: Extract 10,000 uniform surface points with k-NN normal vectors (.PLY).",
        "Interactive 60 FPS WebGL Viewer: Build responsive browser canvas (Next.js 15 + R3F) with orbit controls."
    ]

    for i, obj in enumerate(objs):
        obox = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5 + i*0.88), Inches(7.0), Inches(0.78))
        obox.fill.solid()
        obox.fill.fore_color.rgb = CARD_BG
        obox.line.color.rgb = EMERALD_ACCENT
        tf = obox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✓ {obj.split(':')[0]}: {obj.split(':')[1]}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

    if os.path.exists(img_obj):
        s3.shapes.add_picture(img_obj, Inches(8.1), Inches(1.5), height=Inches(4.3))

    val_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    val_box.line.color.rgb = CYAN_ACCENT
    tf_v = val_box.text_frame
    p = tf_v.paragraphs[0]
    p.text = "System Validation: All primary functional objectives were 100% met and verified across Master Acceptance Testing."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 4: SCOPE AND CONSTRAINTS WITH IMAGE
    # -------------------------------------------------------------------------
    s4 = add_blank_slide()
    add_header(s4, "4. Scope & Constraints — Boundaries & Technical Limits", 4)

    scope_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.4), Inches(4.3))
    scope_box.fill.solid()
    scope_box.fill.fore_color.rgb = CARD_BG
    scope_box.line.color.rgb = CYAN_ACCENT
    tf_s = scope_box.text_frame
    tf_s.word_wrap = True
    p = tf_s.paragraphs[0]
    p.text = "System Scope"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    items = [
        "Single RGB photograph (JPG, PNG, WEBP, BMP) up to 25 MB.",
        "Downloadable textured .GLB mesh and .PLY point cloud.",
        "Targeted at e-commerce, AR/VR spatial apps, and gaming."
    ]
    for it in items:
        p2 = tf_s.add_paragraph()
        p2.text = f"• {it}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    const_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.4), Inches(1.5), Inches(3.4), Inches(4.3))
    const_box.fill.solid()
    const_box.fill.fore_color.rgb = CARD_BG
    const_box.line.color.rgb = AMBER_ACCENT
    tf_c = const_box.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "Operational Constraints"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    items_c = [
        "Optimized for single prominent foreground objects.",
        "Executes within 16 GB VRAM GPU ceiling (NVIDIA T4).",
        "Synthesizes occluded rear textures via triplane diffusion."
    ]
    for it in items_c:
        p2 = tf_c.add_paragraph()
        p2.text = f"• {it}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    if os.path.exists(img_scope):
        s4.shapes.add_picture(img_scope, Inches(8.0), Inches(1.5), height=Inches(4.3))

    hw_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    hw_box.fill.solid()
    hw_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    hw_box.line.color.rgb = INDIGO_ACCENT
    tf_hw = hw_box.text_frame
    p = tf_hw.paragraphs[0]
    p.text = "Hardware Environment: NVIDIA T4 GPU (16 GB VRAM), CUDA 12.1, 32 GB System RAM, Windows/Linux OS."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INDIGO_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 5: BLOCK DIAGRAM
    # -------------------------------------------------------------------------
    s5 = add_blank_slide()
    add_header(s5, "5. Block Diagram — System Architecture & Data Pipelines", 5)

    code_box5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(4.3))
    code_box5.fill.solid()
    code_box5.fill.fore_color.rgb = RGBColor(2, 6, 23)
    code_box5.line.color.rgb = INDIGO_ACCENT
    tf_code = code_box5.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = Inches(0.4)
    tf_code.margin_top = Inches(0.3)

    diagram_text = (
        "+-----------------------------------------------------------------------------------+\n"
        "|                                 CLIENT VIEWPORT                                   |\n"
        "|   +---------------------------------------------------------------------------+   |\n"
        "|   | Next.js 15 Web Dashboard (React 19 / React Three Fiber 60 FPS Canvas)     |   |\n"
        "|   +---------------------------------------------------------------------------+   |\n"
        "+------------------------------------------+----------------------------------------+\n"
        "                                           | HTTP / REST API (JWT Auth)\n"
        "                                           v\n"
        "+-----------------------------------------------------------------------------------+\n"
        "|                             FASTAPI ASYNC BACKEND SERVER                          |\n"
        "|   +---------------------------------------------------------------------------+   |\n"
        "|   | Job Orchestrator & Storage Manager (PyTorch / CUDA 12.1 Engine)           |   |\n"
        "|   +---------------------------------------------------------------------------+   |\n"
        "+------------------------------------------+----------------------------------------+\n"
        "                                           | Sequential CUDA Execution\n"
        "                                           v\n"
        "+-----------------------------------------------------------------------------------+\n"
        "|                             GENERATIVE AI MODEL PIPELINE                          |\n"
        "|  [ Stage 1: SAM 2.1 ] ----> [ Stage 2: Hunyuan3D-2 ] ----> [ Stage 3: Open3D ]    |\n"
        "|  Foreground Cutout           Textured .GLB Mesh             10k .PLY Point Cloud  |\n"
        "+------------------------------------------+----------------------------------------+\n"
        "                                           | Persistent Storage\n"
        "                                           v\n"
        "+-----------------------------------------------------------------------------------+\n"
        "|                         SUPABASE CLOUD DATABASE & STORAGE                         |\n"
        "|   PostgreSQL Jobs Table  |  Row Level Security (RLS)  |  Object Storage Buckets   |\n"
        "+-----------------------------------------------------------------------------------+"
    )
    p = tf_code.paragraphs[0]
    p.text = diagram_text
    p.font.name = 'Courier New'
    p.font.size = Pt(8.5)
    p.font.color.rgb = INDIGO_ACCENT

    b_text = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    b_text.fill.solid()
    b_text.fill.fore_color.rgb = RGBColor(15, 23, 42)
    b_text.line.color.rgb = CYAN_ACCENT
    tf_bt = b_text.text_frame
    p = tf_bt.paragraphs[0]
    p.text = "Pipeline Strategy: Sequential execution & CUDA cache clearing (torch.cuda.empty_cache()) prevent out-of-memory errors."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 6: REQUIREMENTS
    # -------------------------------------------------------------------------
    s6 = add_blank_slide()
    add_header(s6, "6. Requirements — Functional & Non-Functional Specifications", 6)

    fr_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.733), Inches(5.4))
    fr_box.fill.solid()
    fr_box.fill.fore_color.rgb = CARD_BG
    fr_box.line.color.rgb = CYAN_ACCENT
    tf_fr = fr_box.text_frame
    tf_fr.word_wrap = True
    p = tf_fr.paragraphs[0]
    p.text = "6.1 Functional Requirements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    fr_items = [
        ("FR-01 (Authentication)", "Supabase JWT user login & session control."),
        ("FR-02 (Image Ingestion)", "Drag-and-drop file upload with format check & 25 MB cap."),
        ("FR-03 (Segmentation)", "Automated SAM 2.1 background removal & alpha matting."),
        ("FR-04 (Mesh Synthesis)", "Rapid textured 3D mesh synthesis (.GLB format)."),
        ("FR-05 (Point Cloud)", "Poisson disk point cloud sampling with normals (.PLY)."),
        ("FR-06 (3D Viewport)", "Interactive 60 FPS WebGL canvas with orbit controls.")
    ]
    for title, desc in fr_items:
        p2 = tf_fr.add_paragraph()
        p2.text = f"• {title}: {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    nfr_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.733), Inches(5.4))
    nfr_box.fill.solid()
    nfr_box.fill.fore_color.rgb = CARD_BG
    nfr_box.line.color.rgb = INDIGO_ACCENT
    tf_nfr = nfr_box.text_frame
    tf_nfr.word_wrap = True
    p = tf_nfr.paragraphs[0]
    p.text = "6.2 Non-Functional Requirements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = INDIGO_ACCENT

    nfr_items = [
        ("Performance", "Total pipeline latency completed in ~19.6 seconds."),
        ("Memory Ceiling", "Peak GPU memory strictly capped at 11.4 GB VRAM (< 16 GB)."),
        ("Usability & Aesthetics", "Modern glassmorphism UI supporting Light/Dark mode."),
        ("Reliability & Testing", "100% test pass rate across Master Acceptance Matrix (TC-01 to TC-08).")
    ]
    for title, desc in nfr_items:
        p2 = tf_nfr.add_paragraph()
        p2.text = f"• {title}: {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(14)

    # -------------------------------------------------------------------------
    # SLIDE 7: USE-CASE DIAGRAM
    # -------------------------------------------------------------------------
    s7 = add_blank_slide()
    add_header(s7, "7. Use-case Diagram — User & Backend System Interaction", 7)

    code_box7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(4.3))
    code_box7.fill.solid()
    code_box7.fill.fore_color.rgb = RGBColor(2, 6, 23)
    code_box7.line.color.rgb = ROSE_ACCENT
    tf_c7 = code_box7.text_frame
    tf_c7.word_wrap = True
    tf_c7.margin_left = Inches(0.4)
    tf_c7.margin_top = Inches(0.3)

    uc_diagram_text = (
        "                              +-------------------------------------------+\n"
        "                              |         SINGLE-IMAGE 3D SYSTEM            |\n"
        "                              |                                           |\n"
        "    +--------------+          |  (UC-01: User Login / Register)           |          +------------------+\n"
        "    |              | -------- |                                           | -------- |                  |\n"
        "    |              | -------- |  (UC-02: Upload Single RGB Image)         | -------- |                  |\n"
        "    |    USER      |          |                                           |          |  FASTAPI BACKEND |\n"
        "    |   (Client)   | -------- |  (UC-03: Track Processing Progress)       | -------- |    & AI MODELS   |\n"
        "    |              |          |                                           |          |                  |\n"
        "    |              | -------- |  (UC-04: Interact with WebGL 3D Canvas)   | -------- |                  |\n"
        "    |              | -------- |                                           |          |                  |\n"
        "    +--------------+          |  (UC-05: Download GLB / PLY Assets)       |          +------------------+\n"
        "                              |                                           |\n"
        "                              +-------------------------------------------+"
    )
    p = tf_c7.paragraphs[0]
    p.text = uc_diagram_text
    p.font.name = 'Courier New'
    p.font.size = Pt(9.5)
    p.font.color.rgb = ROSE_ACCENT

    actors_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    actors_box.fill.solid()
    actors_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    actors_box.line.color.rgb = CYAN_ACCENT
    tf_act = actors_box.text_frame
    p = tf_act.paragraphs[0]
    p.text = "Actors: User (Uploads, tracks, interacts with 3D canvas, downloads) | Backend System (Validates JWT, executes models, persists data)."
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 8: CONCLUSION AND FUTURE SCOPE
    # -------------------------------------------------------------------------
    s8 = add_blank_slide()
    add_header(s8, "8. Conclusion and Future Scope — Summary & Roadmap", 8)

    conc_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.733), Inches(4.3))
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = CARD_BG
    conc_box.line.color.rgb = AMBER_ACCENT
    tf_conc = conc_box.text_frame
    tf_conc.word_wrap = True
    p = tf_conc.paragraphs[0]
    p.text = "Project Conclusion"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    c_bullets = [
        "Successfully engineered an automated single-image 2D-to-3D asset and point cloud generation web system.",
        "Achieved fast end-to-end execution (~19.6s) and low GPU memory footprint (11.4 GB VRAM) on standard workstation hardware.",
        "Delivered production-ready textured .GLB polygon meshes (~24.5k vertices) and 10,000-point Poisson-sampled .PLY point clouds."
    ]
    for b in c_bullets:
        p2 = tf_conc.add_paragraph()
        p2.text = f"• {b}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(10)

    fut_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.733), Inches(4.3))
    fut_box.fill.solid()
    fut_box.fill.fore_color.rgb = CARD_BG
    fut_box.line.color.rgb = CYAN_ACCENT
    tf_fut = fut_box.text_frame
    tf_fut.word_wrap = True
    p = tf_fut.paragraphs[0]
    p.text = "Future Scope & Enhancements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    f_bullets = [
        "1. PBR Material Texture Synthesis: Generate metallic, roughness, and normal maps for AAA game engine materials.",
        "2. Real-Time WebGPU Gaussian Splatting: Implement 3D Gaussian Splatting for photorealistic volumetric radiance fields.",
        "3. CAD Topology Optimization & 3D Printing: Quad-mesh remashing for watertight .STL additive manufacturing export.",
        "4. WebXR Mobile AR/VR Passthrough: Enable real-world AR asset placement via smartphone cameras and Apple Vision Pro."
    ]
    for b in f_bullets:
        p2 = tf_fut.add_paragraph()
        p2.text = f"• {b}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    stat_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    stat_box.line.color.rgb = EMERALD_ACCENT
    tf_st = stat_box.text_frame
    p = tf_st.paragraphs[0]
    p.text = "Project Status: Complete MCA IV Semester Main Project successfully validated and ready for production deployment."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 9: REFERENCES
    # -------------------------------------------------------------------------
    s9 = add_blank_slide()
    add_header(s9, "9. References — Academic & Literature Citations", 9)

    ref_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = CARD_BG
    ref_box.line.color.rgb = CYAN_ACCENT
    tf_ref = ref_box.text_frame
    tf_ref.word_wrap = True
    p = tf_ref.paragraphs[0]
    p.text = "Academic Citations & Literature"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    refs = [
        "1. Tencent Hunyuan3D-2 Team. (2025). Hunyuan3D 2.0: Scaling Diffusion Models for High-Fidelity 3D Asset Generation. arXiv:2501.12211.",
        "2. Kirillov, A., et al. (2023). Segment Anything. Proceedings of the IEEE/CVF International Conference on Computer Vision (ICCV), pp. 4015-4026.",
        "3. Ravi, N., et al. (2024). SAM 2: Segment Anything in Images and Videos. arXiv:2408.00714.",
        "4. Li, J., et al. (2022). BLIP: Bootstrapping Language-Image Pre-training. ICML, PMLR, 12888-12900.",
        "5. Zhou, Q., et al. (2018). Open3D: A Modern Library for 3D Data Processing. arXiv:1801.09847.",
        "6. Lorensen, W. E., & Cline, H. E. (1987). Marching Cubes: A 3D Surface Construction Algorithm. ACM SIGGRAPH, 21(4), 163-169.",
        "7. Next.js & React Three Fiber Documentation. (2025). Vercel & Poimandres Framework Specifications."
    ]

    for r in refs:
        p2 = tf_ref.add_paragraph()
        p2.text = r
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # Save presentation
    output_path = os.path.abspath(r"c:\Personal\3D\docs\presentation_slides.pptx")
    prs.save(output_path)
    print(f"SUCCESS: Created PPTX file with embedded images at {output_path}")

if __name__ == "__main__":
    create_presentation()
