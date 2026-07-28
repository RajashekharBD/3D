import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Modern Color Palette
    BG_COLOR = RGBColor(11, 15, 25)         # #0B0F19 Midnight Navy
    CARD_BG = RGBColor(15, 23, 42)          # #0F172A Dark Slate Glass
    CARD_BORDER = RGBColor(56, 189, 248)    # #38BDF8 Glowing Cyan
    TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC Pure White
    TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8 Muted Grey-Blue
    CYAN_ACCENT = RGBColor(56, 189, 248)   # #38BDF8 Electric Cyan
    INDIGO_ACCENT = RGBColor(129, 140, 248) # #818CF8 Soft Violet
    EMERALD_ACCENT = RGBColor(52, 211, 153) # #34D399 Neon Emerald
    AMBER_ACCENT = RGBColor(251, 191, 36)   # #FBBF24 Amber Glow
    ROSE_ACCENT = RGBColor(251, 113, 133)   # #FB7185 Rose Fuchsia

    logo_path = os.path.abspath(r"c:\Personal\3D\docs\kle_tech_logo.png")
    pdf_assets_dir = os.path.abspath(r"c:\Personal\3D\docs\pdf_extracted_assets")

    img_cover = os.path.join(pdf_assets_dir, "img_page_1_1.png")
    img_intro = os.path.join(pdf_assets_dir, "img_page_3_1.png")
    img_prob = os.path.join(pdf_assets_dir, "img_page_4_1.png")
    img_obj = os.path.join(pdf_assets_dir, "img_page_5_1.png")
    img_scope = os.path.join(pdf_assets_dir, "img_page_6_1.png")

    def add_blank_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        return slide

    def add_header(slide, title_text, slide_num):
        # Header Container Shape
        header_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.35), Inches(12.133), Inches(0.85))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = CARD_BG
        header_shape.line.color.rgb = RGBColor(30, 41, 59)
        
        # Logo on Header
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.43), height=Inches(0.68))
        
        # Header Text Box
        tx_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.38), Inches(8.8), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "KLE Technological University — Department of MCA"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = CYAN_ACCENT

        # Slide Number Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(0.52), Inches(1.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
        badge.line.color.rgb = CYAN_ACCENT
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = f"Slide {slide_num} / 10"
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 0: TITLE COVER SLIDE (PERFECT 2-COLUMN ALIGNMENT)
    # -------------------------------------------------------------------------
    s0 = add_blank_slide()
    
    # Outer Background Card
    card0 = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.5), Inches(12.133), Inches(6.5))
    card0.fill.solid()
    card0.fill.fore_color.rgb = CARD_BG
    card0.line.color.rgb = CARD_BORDER

    # College Logo Top-Right
    if os.path.exists(logo_path):
        s0.shapes.add_picture(logo_path, Inches(10.7), Inches(0.8), height=Inches(0.9))

    # Left Text Box (Width: 6.5" - Absolutely NO overlap with right image at 7.8")
    t_box0 = s0.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(6.5), Inches(3.4))
    tf0 = t_box0.text_frame
    tf0.word_wrap = True

    p0 = tf0.paragraphs[0]
    p0.text = "MCA IV SEMESTER MAIN PROJECT PRESENTATION"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = CYAN_ACCENT

    p1 = tf0.add_paragraph()
    p1.text = "Automated Single-Image to 3D Asset &\nPoint Cloud Generation System"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_before = Pt(8)

    p2 = tf0.add_paragraph()
    p2.text = "An end-to-end web system converting single 2D RGB photographs into textured 3D polygon meshes (.GLB) and surface point clouds (.PLY) in ~19.6 seconds."
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)

    # Submitted By Card
    sub_box = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.5), Inches(3.1), Inches(1.8))
    sub_box.fill.solid()
    sub_box.fill.fore_color.rgb = RGBColor(2, 6, 23)
    sub_box.line.color.rgb = RGBColor(30, 41, 59)
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = Inches(0.2)
    p = tf_sub.paragraphs[0]
    p.text = "SUBMITTED BY"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p = tf_sub.add_paragraph()
    p.text = "Rajashekhar B Durgad"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_sub.add_paragraph()
    p.text = "USN: 01FE24MCA027"
    p.font.size = Pt(10)
    p.font.color.rgb = CYAN_ACCENT

    # Guidance Card
    guid_box = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(4.5), Inches(3.2), Inches(1.8))
    guid_box.fill.solid()
    guid_box.fill.fore_color.rgb = RGBColor(2, 6, 23)
    guid_box.line.color.rgb = RGBColor(30, 41, 59)
    tf_guid = guid_box.text_frame
    tf_guid.word_wrap = True
    tf_guid.margin_left = Inches(0.2)
    p = tf_guid.paragraphs[0]
    p.text = "UNDER GUIDANCE OF"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p = tf_guid.add_paragraph()
    p.text = "Prof. Akash Hulkund"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p = tf_guid.add_paragraph()
    p.text = "Dept. of MCA, KLE Tech"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # Right Column Image Card (Starts at 7.8" - Completely separated!)
    img_card0 = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.4), Inches(4.5))
    img_card0.fill.solid()
    img_card0.fill.fore_color.rgb = RGBColor(2, 6, 23)
    img_card0.line.color.rgb = CYAN_ACCENT

    if os.path.exists(img_cover):
        s0.shapes.add_picture(img_cover, Inches(8.0), Inches(2.0), width=Inches(4.0), height=Inches(4.1))

    # -------------------------------------------------------------------------
    # SLIDE 1: INTRODUCTION
    # -------------------------------------------------------------------------
    s1 = add_blank_slide()
    add_header(s1, "1. Introduction — Overview & Multi-Stage Architecture", 1)

    cards = [
        ("Automated 2D to 3D Pipeline", "Transforms raw 2D photographs into textured 3D polygon meshes (.GLB) and surface point clouds (.PLY) without multi-view setups or manual modeling.", CYAN_ACCENT),
        ("Multi-Stage AI Engine", "Integrates Meta SAM 2.1 (foreground isolation), Tencent Hunyuan3D-2 (triplane synthesis), and Open3D (Poisson disk surface sampling).", INDIGO_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(cards):
        c_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4 + i*2.2), Inches(7.0), Inches(2.0))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = CARD_BG
        c_box.line.color.rgb = color
        tf = c_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # Right Image Card
    img_card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.4), Inches(4.8), Inches(4.2))
    img_card1.fill.solid()
    img_card1.fill.fore_color.rgb = RGBColor(2, 6, 23)
    img_card1.line.color.rgb = CYAN_ACCENT

    if os.path.exists(img_intro):
        s1.shapes.add_picture(img_intro, Inches(8.1), Inches(1.55), width=Inches(4.4), height=Inches(3.9))

    hl = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    hl.fill.solid()
    hl.fill.fore_color.rgb = RGBColor(15, 23, 42)
    hl.line.color.rgb = EMERALD_ACCENT
    tf_hl = hl.text_frame
    p = tf_hl.paragraphs[0]
    p.text = "Key Highlight: Complete processing finished in ~19.6 seconds with a peak GPU memory footprint of 11.4 GB VRAM."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT
    # -------------------------------------------------------------------------
    s2 = add_blank_slide()
    add_header(s2, "2. Problem Statement — Challenges & Technical Bottlenecks", 2)

    probs = [
        ("High Manual Overhead in 3D Modeling", "Manual CAD software (Blender, Maya) requires specialized artistic skills and hours to days of labor per single 3D asset.", ROSE_ACCENT),
        ("Photogrammetry & Text-to-3D Limits", "Photogrammetry requires 50+ calibrated camera shots. Text-to-3D generative models lack spatial precision for exact physical objects.", AMBER_ACCENT),
        ("GPU Hardware & VRAM Bottlenecks", "Simultaneous execution of multiple heavy neural networks causes Out-Of-Memory (OOM) crashes on standard workstation GPUs.", INDIGO_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(probs):
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4 + i*1.45), Inches(7.0), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)

    # Right Image Card
    img_card2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.4), Inches(4.8), Inches(4.2))
    img_card2.fill.solid()
    img_card2.fill.fore_color.rgb = RGBColor(2, 6, 23)
    img_card2.line.color.rgb = ROSE_ACCENT

    if os.path.exists(img_prob):
        s2.shapes.add_picture(img_prob, Inches(8.1), Inches(1.55), width=Inches(4.4), height=Inches(3.9))

    g_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    g_box.line.color.rgb = CYAN_ACCENT
    tf_g = g_box.text_frame
    p = tf_g.paragraphs[0]
    p.text = "The Core Goal: Provide a prompt-free, single-image 3D synthesis web system operating reliably under a 16 GB GPU VRAM ceiling."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 3: OBJECTIVES
    # -------------------------------------------------------------------------
    s3 = add_blank_slide()
    add_header(s3, "3. Objectives — System Technical Objectives", 3)

    objs = [
        "Prompt-Free 2D-to-3D Pipeline: Automate single RGB image ingestion without text prompts.",
        "Zero-Shot Background Isolation: Strip background clutter using Meta SAM 2.1 for clean RGBA cutouts.",
        "Rapid 3D Mesh Synthesis: Generate textured watertight .GLB polygon meshes in ~15s using Hunyuan3D-2.",
        "Point Cloud Extraction: Extract 10,000 uniform surface points with k-NN normal vectors (.PLY).",
        "Interactive 60 FPS WebGL Viewer: Build responsive browser canvas (Next.js 15 + R3F) with orbit controls."
    ]

    for i, obj in enumerate(objs):
        obox = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4 + i*0.88), Inches(7.0), Inches(0.78))
        obox.fill.solid()
        obox.fill.fore_color.rgb = CARD_BG
        obox.line.color.rgb = EMERALD_ACCENT
        tf = obox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = f"✓ {obj.split(':')[0]}: {obj.split(':')[1]}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

    # Right Image Card
    img_card3 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.4), Inches(4.8), Inches(4.2))
    img_card3.fill.solid()
    img_card3.fill.fore_color.rgb = RGBColor(2, 6, 23)
    img_card3.line.color.rgb = EMERALD_ACCENT

    if os.path.exists(img_obj):
        s3.shapes.add_picture(img_obj, Inches(8.1), Inches(1.55), width=Inches(4.4), height=Inches(3.9))

    val_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    val_box.line.color.rgb = CYAN_ACCENT
    tf_v = val_box.text_frame
    p = tf_v.paragraphs[0]
    p.text = "System Validation: All primary functional objectives were 100% met and verified across Master Acceptance Testing."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 4: SCOPE AND CONSTRAINTS
    # -------------------------------------------------------------------------
    s4 = add_blank_slide()
    add_header(s4, "4. Scope & Constraints — Boundaries & Technical Limits", 4)

    scope_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(3.4), Inches(4.3))
    scope_box.fill.solid()
    scope_box.fill.fore_color.rgb = CARD_BG
    scope_box.line.color.rgb = CYAN_ACCENT
    tf_s = scope_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = Inches(0.2)
    p = tf_s.paragraphs[0]
    p.text = "System Scope"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    items = [
        "Single RGB photograph (JPG, PNG, WEBP, BMP) up to 25 MB.",
        "Downloadable textured .GLB mesh and .PLY point cloud.",
        "Targeted at e-commerce, AR/VR spatial apps, and gaming."
    ]
    for it in items:
        p2 = tf_s.add_paragraph()
        p2.text = f"• {it}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    const_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.4), Inches(3.4), Inches(4.3))
    const_box.fill.solid()
    const_box.fill.fore_color.rgb = CARD_BG
    const_box.line.color.rgb = AMBER_ACCENT
    tf_c = const_box.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.2)
    p = tf_c.paragraphs[0]
    p.text = "Operational Constraints"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    items_c = [
        "Optimized for single prominent foreground objects.",
        "Executes within 16 GB VRAM GPU ceiling (NVIDIA T4).",
        "Synthesizes occluded rear textures via triplane diffusion."
    ]
    for it in items_c:
        p2 = tf_c.add_paragraph()
        p2.text = f"• {it}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # Right Image Card
    img_card4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.4), Inches(4.933), Inches(4.3))
    img_card4.fill.solid()
    img_card4.fill.fore_color.rgb = RGBColor(2, 6, 23)
    img_card4.line.color.rgb = INDIGO_ACCENT

    if os.path.exists(img_scope):
        s4.shapes.add_picture(img_scope, Inches(8.0), Inches(1.55), width=Inches(4.5), height=Inches(4.0))

    hw_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    hw_box.fill.solid()
    hw_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    hw_box.line.color.rgb = INDIGO_ACCENT
    tf_hw = hw_box.text_frame
    p = tf_hw.paragraphs[0]
    p.text = "Hardware Environment: NVIDIA T4 GPU (16 GB VRAM), CUDA 12.1, 32 GB System RAM, Windows/Linux OS."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INDIGO_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 5: BLOCK DIAGRAM (HIGH-RES GRAPHICAL DIAGRAM)
    # -------------------------------------------------------------------------
    s5 = add_blank_slide()
    add_header(s5, "5. Block Diagram — System Architecture & Data Pipelines", 5)

    arch_img_path = os.path.abspath(r"c:\Personal\3D\docs\system_architecture_diagram.png")
    if os.path.exists(arch_img_path):
        card5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.35), Inches(12.133), Inches(4.4))
        card5.fill.solid()
        card5.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card5.line.color.rgb = CYAN_ACCENT
        s5.shapes.add_picture(arch_img_path, Inches(0.8), Inches(1.45), width=Inches(11.733), height=Inches(4.2))

    b_text = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    b_text.fill.solid()
    b_text.fill.fore_color.rgb = RGBColor(15, 23, 42)
    b_text.line.color.rgb = CYAN_ACCENT
    tf_bt = b_text.text_frame
    p = tf_bt.paragraphs[0]
    p.text = "Pipeline Strategy: Sequential CUDA execution & GPU cache clearing (torch.cuda.empty_cache()) prevent VRAM OOM crashes."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 6: REQUIREMENTS
    # -------------------------------------------------------------------------
    s6 = add_blank_slide()
    add_header(s6, "6. Requirements — Functional & Non-Functional Specifications", 6)

    fr_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.4))
    fr_box.fill.solid()
    fr_box.fill.fore_color.rgb = CARD_BG
    fr_box.line.color.rgb = CYAN_ACCENT
    tf_fr = fr_box.text_frame
    tf_fr.word_wrap = True
    tf_fr.margin_left = Inches(0.3)
    p = tf_fr.paragraphs[0]
    p.text = "6.1 Functional Requirements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    fr_items = [
        ("FR-01 (Authentication)", "Supabase JWT user login & session control."),
        ("FR-02 (Image Ingestion)", "Drag-and-drop file upload with format check & 25 MB cap."),
        ("FR-03 (Segmentation)", "Automated SAM 2.1 background removal & alpha matting."),
        ("FR-04 (Mesh Synthesis)", "Rapid textured 3D mesh synthesis (.GLB format)."),
        ("FR-05 (Point Cloud)", "Poisson disk point cloud sampling with normals (.PLY)."),
        ("FR-06 (3D Viewport)", "Interactive 60 FPS WebGL canvas with orbit controls.")
    ]
    for title, desc in fr_items:
        p2 = tf_fr.add_paragraph()
        p2.text = f"• {title}: {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    nfr_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.9), Inches(5.4))
    nfr_box.fill.solid()
    nfr_box.fill.fore_color.rgb = CARD_BG
    nfr_box.line.color.rgb = INDIGO_ACCENT
    tf_nfr = nfr_box.text_frame
    tf_nfr.word_wrap = True
    tf_nfr.margin_left = Inches(0.3)
    p = tf_nfr.paragraphs[0]
    p.text = "6.2 Non-Functional Requirements"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = INDIGO_ACCENT

    nfr_items = [
        ("Performance", "Total pipeline latency completed in ~19.6 seconds."),
        ("Memory Ceiling", "Peak GPU memory strictly capped at 11.4 GB VRAM (< 16 GB)."),
        ("Usability & Aesthetics", "Modern glassmorphism UI supporting Light/Dark mode."),
        ("Reliability & Testing", "100% test pass rate across Master Acceptance Matrix (TC-01 to TC-08).")
    ]
    for title, desc in nfr_items:
        p2 = tf_nfr.add_paragraph()
        p2.text = f"• {title}: {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(14)

    # -------------------------------------------------------------------------
    # SLIDE 7: USE-CASE DIAGRAM (HIGH-RES GRAPHICAL DIAGRAM)
    # -------------------------------------------------------------------------
    s7 = add_blank_slide()
    add_header(s7, "7. Use-case Diagram — User & Backend System Interaction", 7)

    uc_img_path = os.path.abspath(r"c:\Personal\3D\docs\usecase_diagram_graphic.png")
    if os.path.exists(uc_img_path):
        card7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.35), Inches(12.133), Inches(4.4))
        card7.fill.solid()
        card7.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card7.line.color.rgb = ROSE_ACCENT
        s7.shapes.add_picture(uc_img_path, Inches(1.8), Inches(1.45), width=Inches(9.733), height=Inches(4.2))

    actors_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    actors_box.fill.solid()
    actors_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    actors_box.line.color.rgb = CYAN_ACCENT
    tf_act = actors_box.text_frame
    p = tf_act.paragraphs[0]
    p.text = "Actors: User (Uploads, tracks, interacts with 3D canvas, downloads) | Backend System (Validates JWT, executes models, persists data)."
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 8: CONCLUSION AND FUTURE SCOPE
    # -------------------------------------------------------------------------
    s8 = add_blank_slide()
    add_header(s8, "8. Conclusion and Future Scope — Summary & Roadmap", 8)

    conc_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.9), Inches(4.3))
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = CARD_BG
    conc_box.line.color.rgb = AMBER_ACCENT
    tf_conc = conc_box.text_frame
    tf_conc.word_wrap = True
    tf_conc.margin_left = Inches(0.3)
    p = tf_conc.paragraphs[0]
    p.text = "Project Conclusion"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    c_bullets = [
        "Successfully engineered an automated single-image 2D-to-3D asset and point cloud generation web system.",
        "Achieved fast end-to-end execution (~19.6s) and low GPU memory footprint (11.4 GB VRAM) on standard workstation hardware.",
        "Delivered production-ready textured .GLB polygon meshes (~24.5k vertices) and 10,000-point Poisson-sampled .PLY point clouds."
    ]
    for b in c_bullets:
        p2 = tf_conc.add_paragraph()
        p2.text = f"• {b}"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    fut_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.9), Inches(4.3))
    fut_box.fill.solid()
    fut_box.fill.fore_color.rgb = CARD_BG
    fut_box.line.color.rgb = CYAN_ACCENT
    tf_fut = fut_box.text_frame
    tf_fut.word_wrap = True
    tf_fut.margin_left = Inches(0.3)
    p = tf_fut.paragraphs[0]
    p.text = "Future Scope & Enhancements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    f_bullets = [
        "1. PBR Material Texture Synthesis: Generate metallic, roughness, and normal maps for AAA game engine materials.",
        "2. Real-Time WebGPU Gaussian Splatting: Implement 3D Gaussian Splatting for photorealistic volumetric radiance fields.",
        "3. CAD Topology Optimization & 3D Printing: Quad-mesh remashing for watertight .STL additive manufacturing export.",
        "4. WebXR Mobile AR/VR Passthrough: Enable real-world AR asset placement via smartphone cameras and Apple Vision Pro."
    ]
    for b in f_bullets:
        p2 = tf_fut.add_paragraph()
        p2.text = f"• {b}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    stat_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.133), Inches(0.9))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    stat_box.line.color.rgb = EMERALD_ACCENT
    tf_st = stat_box.text_frame
    p = tf_st.paragraphs[0]
    p.text = "Project Status: Complete MCA IV Semester Main Project successfully validated and ready for production deployment."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 9: REFERENCES
    # -------------------------------------------------------------------------
    s9 = add_blank_slide()
    add_header(s9, "9. References — Academic & Literature Citations", 9)

    ref_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = CARD_BG
    ref_box.line.color.rgb = CYAN_ACCENT
    tf_ref = ref_box.text_frame
    tf_ref.word_wrap = True
    tf_ref.margin_left = Inches(0.3)
    p = tf_ref.paragraphs[0]
    p.text = "Academic Citations & Literature"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    refs = [
        "1. Tencent Hunyuan3D-2 Team. (2025). Hunyuan3D 2.0: Scaling Diffusion Models for High-Fidelity 3D Asset Generation. arXiv:2501.12211.",
        "2. Kirillov, A., et al. (2023). Segment Anything. Proceedings of the IEEE/CVF International Conference on Computer Vision (ICCV), pp. 4015-4026.",
        "3. Ravi, N., et al. (2024). SAM 2: Segment Anything in Images and Videos. arXiv:2408.00714.",
        "4. Li, J., et al. (2022). BLIP: Bootstrapping Language-Image Pre-training. ICML, PMLR, 12888-12900.",
        "5. Zhou, Q., et al. (2018). Open3D: A Modern Library for 3D Data Processing. arXiv:1801.09847.",
        "6. Lorensen, W. E., & Cline, H. E. (1987). Marching Cubes: A 3D Surface Construction Algorithm. ACM SIGGRAPH, 21(4), 163-169.",
        "7. Next.js & React Three Fiber Documentation. (2025). Vercel & Poimandres Framework Specifications."
    ]

    for r in refs:
        p2 = tf_ref.add_paragraph()
        p2.text = r
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # Save presentation
    output_path = os.path.abspath(r"c:\Personal\3D\docs\presentation_slides_v4.pptx")
    prs.save(output_path)
    print(f"PERFECT GRAPHICAL USE CASE DIAGRAM SUCCESS: Created PPTX file at {output_path}")

    # Try saving to main filenames if not locked
    for p_name in [r"c:\Personal\3D\docs\presentation_slides.pptx", r"c:\Personal\3D\docs\presentation_slides_v2.pptx", r"c:\Personal\3D\docs\presentation_slides_v3.pptx"]:
        try:
            prs.save(os.path.abspath(p_name))
            print(f"Updated {p_name} as well!")
        except Exception:
            pass

if __name__ == "__main__":
    create_presentation()
